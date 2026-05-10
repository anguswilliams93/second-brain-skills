from pathlib import Path
import math
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_BG = "0d2e28"
BRAND_BG_ALT = "0a231e"
BRAND_TEXT = "ffffff"
BRAND_TEXT_SECONDARY = "a1e8d9"
BRAND_ACCENT = "4dd4c6"
BRAND_ACCENT_SECONDARY = "2dd4aa"
BRAND_ACCENT_TERTIARY = "6ffce0"
BRAND_CARD_BG = "113a33"
BRAND_CARD_BG_ALT = "15473e"
BRAND_HEADING_FONT = "Inter Tight"
BRAND_BODY_FONT = "Inter"


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)
    return s


# === Slide 1: Title ===
slide = new_slide()
top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
top_bar.fill.solid(); top_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); top_bar.line.fill.background()

eyebrow = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.83), Inches(0.5))
p = eyebrow.text_frame.paragraphs[0]
p.text = "METHODOLOGY  \u00B7  PRODUCTS  \u00B7  LOOP  \u00B7  MODEL"
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(14); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY); p.alignment = PP_ALIGN.CENTER

hb = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(11.83), Inches(1.8))
tf = hb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "The BlockScore Livability Index"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(56); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT); p.alignment = PP_ALIGN.CENTER

sb = slide.shapes.add_textbox(Inches(1.5), Inches(4.7), Inches(10.33), Inches(1.0))
p = sb.text_frame.paragraphs[0]
p.text = "Measuring what makes an Australian strata building liveable"
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(24); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY); p.alignment = PP_ALIGN.CENTER

al = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.0), Inches(5.9), Inches(1.333), Inches(0.05))
al.fill.solid(); al.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); al.line.fill.background()

# === Slide 2: Section break 01 ===
slide = new_slide()
accent_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
accent_block.fill.solid(); accent_block.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT); accent_block.line.fill.background()

num = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(3.5), Inches(1.5))
p = num.text_frame.paragraphs[0]; p.text = "01"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(120); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_ACCENT); p.alignment = PP_ALIGN.CENTER

divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(1.5), Inches(0.08), Inches(4.5))
divider.fill.solid(); divider.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); divider.line.fill.background()

tb = slide.shapes.add_textbox(Inches(5.2), Inches(2.6), Inches(7.5), Inches(1.5))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "The Methodology"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

sb2 = slide.shapes.add_textbox(Inches(5.2), Inches(4.2), Inches(7.5), Inches(1.0))
tf = sb2.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Three research-validated layers. One benchmarkable score."
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(20); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

# === Slide 3: Floating cards - Three-layer design ===
slide = new_slide()
tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.text = "A three-layer design"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

CARDS_L = [
    ("A \u00B7 Wellbeing", "How residents feel about life here. Life satisfaction, safety, community, future security \u2014 benchmarked to Australian wellbeing norms."),
    ("B \u00B7 Environment", "What the building and neighbourhood are actually like. Architecture, services, social fit, maintenance, lifestyle."),
    ("C \u00B7 Context", "The daily-life context shaping experience \u2014 safety, information access, leisure and environmental quality."),
]
configs = [
    {"x": 0.8, "y": 2.0, "accent": BRAND_ACCENT, "bg": BRAND_CARD_BG},
    {"x": 4.5, "y": 2.4, "accent": BRAND_ACCENT_SECONDARY, "bg": BRAND_CARD_BG_ALT},
    {"x": 8.2, "y": 1.8, "accent": BRAND_ACCENT_TERTIARY, "bg": BRAND_CARD_BG},
]

for i, (t, d) in enumerate(CARDS_L):
    c = configs[i]
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(c["x"] + 0.15), Inches(c["y"] + 0.15), Inches(4.2), Inches(4.0))
    shadow.fill.solid(); shadow.fill.fore_color.rgb = hex_to_rgb("061713"); shadow.line.fill.background()
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(c["x"]), Inches(c["y"]), Inches(4.2), Inches(4.0))
    card.fill.solid(); card.fill.fore_color.rgb = hex_to_rgb(c["bg"]); card.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(c["x"] + 0.3), Inches(c["y"] + 0.3), Inches(0.8), Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = hex_to_rgb(c["accent"]); bar.line.fill.background()
    nb = slide.shapes.add_textbox(Inches(c["x"] + 0.3), Inches(c["y"] + 0.6), Inches(1.0), Inches(0.8))
    pp = nb.text_frame.paragraphs[0]; pp.text = f"0{i+1}"
    pp.font.name = BRAND_HEADING_FONT; pp.font.size = Pt(32); pp.font.bold = True; pp.font.color.rgb = hex_to_rgb(c["accent"])
    tt = slide.shapes.add_textbox(Inches(c["x"] + 0.3), Inches(c["y"] + 1.5), Inches(3.6), Inches(1.0))
    tff = tt.text_frame; tff.word_wrap = True
    pp = tff.paragraphs[0]; pp.text = t
    pp.font.name = BRAND_HEADING_FONT; pp.font.size = Pt(20); pp.font.bold = True; pp.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    dd = slide.shapes.add_textbox(Inches(c["x"] + 0.3), Inches(c["y"] + 2.5), Inches(3.6), Inches(1.4))
    tff = dd.text_frame; tff.word_wrap = True
    pp = tff.paragraphs[0]; pp.text = d
    pp.font.name = BRAND_BODY_FONT; pp.font.size = Pt(13); pp.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

# === Slide 4: Venn (custom) ===
slide = new_slide()

tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.9))
p = tb.text_frame.paragraphs[0]; p.text = "Livability lives in the overlap"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35), Inches(2.5), Inches(0.06))
underline.fill.solid(); underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); underline.line.fill.background()


def outlined_circle(slide, cx, cy, r, color, stroke_pt=3):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - r), Inches(cy - r), Inches(2 * r), Inches(2 * r))
    shape.fill.background()
    shape.line.color.rgb = hex_to_rgb(color)
    shape.line.width = Pt(stroke_pt)
    return shape


r = 1.7
ax, ay = 6.666, 3.1
bx, by = 5.666, 4.5
cxv, cyv = 7.666, 4.5

outlined_circle(slide, ax, ay, r, BRAND_ACCENT)
outlined_circle(slide, bx, by, r, BRAND_ACCENT_SECONDARY)
outlined_circle(slide, cxv, cyv, r, BRAND_ACCENT_TERTIARY)

# Labels
alab = slide.shapes.add_textbox(Inches(ax - 2.0), Inches(ay - r - 0.55), Inches(4.0), Inches(0.5))
p = alab.text_frame.paragraphs[0]; p.text = "A \u00B7 Subjective Wellbeing"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_ACCENT); p.alignment = PP_ALIGN.CENTER

blab = slide.shapes.add_textbox(Inches(0.3), Inches(by + r - 0.1), Inches(4.5), Inches(0.5))
p = blab.text_frame.paragraphs[0]; p.text = "B \u00B7 Residential Environment"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_ACCENT_SECONDARY); p.alignment = PP_ALIGN.LEFT

clab = slide.shapes.add_textbox(Inches(8.5), Inches(cyv + r - 0.1), Inches(4.5), Inches(0.5))
p = clab.text_frame.paragraphs[0]; p.text = "C \u00B7 Health & QoL Context"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_ACCENT_TERTIARY); p.alignment = PP_ALIGN.RIGHT

# Pairwise intersection hints
pab = slide.shapes.add_textbox(Inches(4.3), Inches(3.8), Inches(1.6), Inches(0.3))
p = pab.text_frame.paragraphs[0]; p.text = "felt experience"
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(9); p.font.italic = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY); p.alignment = PP_ALIGN.CENTER

pac = slide.shapes.add_textbox(Inches(7.5), Inches(3.8), Inches(1.6), Inches(0.3))
p = pac.text_frame.paragraphs[0]; p.text = "personal context"
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(9); p.font.italic = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY); p.alignment = PP_ALIGN.CENTER

pbc = slide.shapes.add_textbox(Inches(5.9), Inches(5.5), Inches(1.6), Inches(0.3))
p = pbc.text_frame.paragraphs[0]; p.text = "contextual quality"
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(9); p.font.italic = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY); p.alignment = PP_ALIGN.CENTER

# BSLI center
bsli_box = slide.shapes.add_textbox(Inches(5.466), Inches(3.75), Inches(2.4), Inches(0.6))
p = bsli_box.text_frame.paragraphs[0]; p.text = "BSLI"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(30); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT); p.alignment = PP_ALIGN.CENTER

bsli_sub = slide.shapes.add_textbox(Inches(5.266), Inches(4.35), Inches(2.8), Inches(0.4))
p = bsli_sub.text_frame.paragraphs[0]; p.text = "Livability"
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(11); p.font.color.rgb = hex_to_rgb(BRAND_ACCENT); p.alignment = PP_ALIGN.CENTER

cap = slide.shapes.add_textbox(Inches(1.0), Inches(6.7), Inches(11.333), Inches(0.6))
p = cap.text_frame.paragraphs[0]
p.text = "The Index only moves when evidence across all three layers does. How they are weighted and combined is BlockScore IP."
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(13); p.font.italic = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY); p.alignment = PP_ALIGN.CENTER

# === Slide 5: Giant focus "0-100" ===
slide = new_slide()
bg_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.5), Inches(0.5), Inches(6.333), Inches(6.333))
bg_circle.fill.solid(); bg_circle.fill.fore_color.rgb = hex_to_rgb("0a231e"); bg_circle.line.fill.background()

ab = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(0.6))
p = ab.text_frame.paragraphs[0]; p.text = "EVERY AUSTRALIAN STRATA BUILDING"
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(14); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY); p.alignment = PP_ALIGN.CENTER

big = slide.shapes.add_textbox(Inches(0), Inches(2.2), Inches(13.333), Inches(3.5))
tf = big.text_frame; tf.word_wrap = False
p = tf.paragraphs[0]; p.text = "0\u2013100"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(200); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_ACCENT); p.alignment = PP_ALIGN.CENTER

bel = slide.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12.333), Inches(0.6))
p = bel.text_frame.paragraphs[0]
p.text = "One benchmarkable Livability score, grounded in peer-reviewed research"
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(20); p.font.color.rgb = hex_to_rgb(BRAND_TEXT); p.alignment = PP_ALIGN.CENTER

aline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(6.4), Inches(2.333), Inches(0.04))
aline.fill.solid(); aline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); aline.line.fill.background()

out = Path("C:/dev/obsidian/second-brain-skills/.claude/skills/pptx-generator/output/blockscore/livability-index-2026-04-18-part1.pptx")
prs.save(out)
print(f"Saved {out}")
