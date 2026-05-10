from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_BG = "0d2e28"
BRAND_BG_ALT = "0a231e"
BRAND_TEXT = "ffffff"
BRAND_TEXT_SECONDARY = "a1e8d9"
BRAND_ACCENT = "4dd4c6"
BRAND_ACCENT_SECONDARY = "2dd4aa"
BRAND_ACCENT_TERTIARY = "6ffce0"
BRAND_CARD_BG = "113a33"
BRAND_CARD_BG_ALT = "15473e"
BRAND_HEADING_FONT = "Inter Tight"
BRAND_BODY_FONT = "Inter"


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)
    return s


# === Slide 16: Floating cards - three funding streams ===
slide = new_slide()
tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1.0))
p = tb.text_frame.paragraphs[0]; p.text = "Three funding streams, one flywheel"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(38); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

CARDS = [
    ("SaaS funds the platform", "Dashboard subscriptions pay for the Livability Index itself \u2014 research, engineering, validation and ML refinement."),
    ("Reports fund the buildings", "Every Building Insight Report sends revenue straight into the featured building\u2019s sinking fund for maintenance."),
    ("Surveys fund the residents", "Residents earn credit against the cost of the report they care about most \u2014 honest feedback, fair reward."),
]
configs = [
    {"x": 0.8, "y": 2.0, "accent": BRAND_ACCENT, "bg": BRAND_CARD_BG},
    {"x": 4.5, "y": 2.4, "accent": BRAND_ACCENT_SECONDARY, "bg": BRAND_CARD_BG_ALT},
    {"x": 8.2, "y": 1.8, "accent": BRAND_ACCENT_TERTIARY, "bg": BRAND_CARD_BG},
]
for i, (t, d) in enumerate(CARDS):
    c = configs[i]
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(c["x"] + 0.15), Inches(c["y"] + 0.15), Inches(4.2), Inches(4.0))
    shadow.fill.solid(); shadow.fill.fore_color.rgb = hex_to_rgb("061713"); shadow.line.fill.background()
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(c["x"]), Inches(c["y"]), Inches(4.2), Inches(4.0))
    card.fill.solid(); card.fill.fore_color.rgb = hex_to_rgb(c["bg"]); card.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(c["x"] + 0.3), Inches(c["y"] + 0.3), Inches(0.8), Inches(0.12))
    bar.fill.solid(); bar.fill.fore_color.rgb = hex_to_rgb(c["accent"]); bar.line.fill.background()
    nb = slide.shapes.add_textbox(Inches(c["x"] + 0.3), Inches(c["y"] + 0.6), Inches(1.0), Inches(0.8))
    pp = nb.text_frame.paragraphs[0]; pp.text = f"0{i + 1}"
    pp.font.name = BRAND_HEADING_FONT; pp.font.size = Pt(32); pp.font.bold = True; pp.font.color.rgb = hex_to_rgb(c["accent"])
    tt = slide.shapes.add_textbox(Inches(c["x"] + 0.3), Inches(c["y"] + 1.5), Inches(3.6), Inches(1.1))
    tff = tt.text_frame; tff.word_wrap = True
    pp = tff.paragraphs[0]; pp.text = t
    pp.font.name = BRAND_HEADING_FONT; pp.font.size = Pt(18); pp.font.bold = True; pp.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    dd = slide.shapes.add_textbox(Inches(c["x"] + 0.3), Inches(c["y"] + 2.6), Inches(3.6), Inches(1.3))
    tff = dd.text_frame; tff.word_wrap = True
    pp = tff.paragraphs[0]; pp.text = d
    pp.font.name = BRAND_BODY_FONT; pp.font.size = Pt(13); pp.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

# === Slide 17: Closing ===
slide = new_slide()

bottom_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.0), Inches(13.333), Inches(2.5))
bottom_block.fill.solid(); bottom_block.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT); bottom_block.line.fill.background()

top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
top_bar.fill.solid(); top_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); top_bar.line.fill.background()

hb = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.33), Inches(1.5))
tf = hb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Know your building. Protect your tenants."
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(52); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT); p.alignment = PP_ALIGN.CENTER

underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.6), Inches(2.333), Inches(0.08))
underline.fill.solid(); underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); underline.line.fill.background()

tag = slide.shapes.add_textbox(Inches(0.5), Inches(3.95), Inches(12.33), Inches(0.5))
p = tag.text_frame.paragraphs[0]
p.text = "Anonymous feedback. Actionable insights. Benchmarkable building scores."
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(18); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY); p.alignment = PP_ALIGN.CENTER

# CTA items
cta_items = [
    "Read the methodology",
    "Join the waitlist",
    "Talk to our team",
]
num_items = len(cta_items)
item_width = 12.0 / num_items
start_x = 0.66

for i, item in enumerate(cta_items):
    x_pos = start_x + (i * item_width)
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + (item_width / 2) - 0.25), Inches(5.3), Inches(0.5), Inches(0.5))
    badge.fill.solid(); badge.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); badge.line.fill.background()

    btext = slide.shapes.add_textbox(Inches(x_pos + (item_width / 2) - 0.25), Inches(5.35), Inches(0.5), Inches(0.5))
    p = btext.text_frame.paragraphs[0]; p.text = str(i + 1)
    p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_BG); p.alignment = PP_ALIGN.CENTER

    ib = slide.shapes.add_textbox(Inches(x_pos), Inches(5.95), Inches(item_width), Inches(0.8))
    tf = ib.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = item
    p.font.name = BRAND_BODY_FONT; p.font.size = Pt(18); p.font.color.rgb = hex_to_rgb(BRAND_TEXT); p.alignment = PP_ALIGN.CENTER

sub = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.33), Inches(0.4))
p = sub.text_frame.paragraphs[0]; p.text = "blockscore.net/methodology"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(14); p.font.color.rgb = hex_to_rgb(BRAND_ACCENT); p.alignment = PP_ALIGN.CENTER

out = Path("C:/dev/obsidian/second-brain-skills/.claude/skills/pptx-generator/output/blockscore/livability-index-2026-04-18-part4.pptx")
prs.save(out)
print(f"Saved {out}")
