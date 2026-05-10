from pptx import Presentation
from pptx.dml.color import RGBColor
from pathlib import Path
import copy
from lxml import etree


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


BRAND_BG = "0d2e28"

output_dir = Path("C:/dev/obsidian/second-brain-skills/.claude/skills/pptx-generator/output/blockscore")
part_files = sorted(output_dir.glob("livability-index-2026-04-18-part*.pptx"))
print(f"Found {len(part_files)} part files")

combined = Presentation(part_files[0])

# Fix backgrounds in already-loaded base (first file slides), just in case
for s in combined.slides:
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

for part_file in part_files[1:]:
    part_prs = Presentation(part_file)
    for src_slide in part_prs.slides:
        blank_layout = combined.slide_layouts[6]
        new_slide = combined.slides.add_slide(blank_layout)
        # CRITICAL: explicit bg
        new_slide.background.fill.solid()
        new_slide.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)

        # Copy shapes
        for shape in src_slide.shapes:
            el = shape.element
            new_el = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

final = output_dir / "BlockScore-Livability-Index-2026-04-18.pptx"
combined.save(final)
print(f"Saved {final}")

# Cleanup parts
for p in part_files:
    p.unlink()
    print(f"Deleted {p.name}")
