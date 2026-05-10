from pathlib import Path
import math
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_BG = "0d2e28"
BRAND_BG_ALT = "0a231e"
BRAND_TEXT = "ffffff"
BRAND_TEXT_SECONDARY = "a1e8d9"
BRAND_ACCENT = "4dd4c6"
BRAND_ACCENT_SECONDARY = "2dd4aa"
BRAND_ACCENT_TERTIARY = "6ffce0"
BRAND_CARD_BG = "113a33"
BRAND_CARD_BG_ALT = "15473e"
BRAND_HEADING_FONT = "Inter Tight"
BRAND_BODY_FONT = "Inter"


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)
    return s


def section_break(number, title, subtitle):
    slide = new_slide()
    accent_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
    accent_block.fill.solid(); accent_block.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT); accent_block.line.fill.background()
    num = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(3.5), Inches(1.5))
    p = num.text_frame.paragraphs[0]; p.text = number
    p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(120); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_ACCENT); p.alignment = PP_ALIGN.CENTER
    divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(1.5), Inches(0.08), Inches(4.5))
    divider.fill.solid(); divider.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); divider.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(5.2), Inches(2.6), Inches(7.5), Inches(1.5))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)
    sb2 = slide.shapes.add_textbox(Inches(5.2), Inches(4.2), Inches(7.5), Inches(1.0))
    tf = sb2.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = subtitle
    p.font.name = BRAND_BODY_FONT; p.font.size = Pt(20); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)


# === Slide 11: Section break 03 The BlockScore Loop ===
section_break("03", "The BlockScore Loop", "Value flows full circle \u2014 from resident to building and back.")

# === Slide 12: Circular-hero Loop (4 nodes: Survey, Credit, Report, Sinking Fund) ===
slide = new_slide()

# Title
tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
p = tb.text_frame.paragraphs[0]; p.text = "The BlockScore Loop"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT); p.alignment = PP_ALIGN.CENTER

center_x = 6.666
center_y = 4.2

# Outer ring
outer_ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - 3.2), Inches(center_y - 3.2), Inches(6.4), Inches(6.4))
outer_ring.fill.background()
outer_ring.line.color.rgb = hex_to_rgb("165046")
outer_ring.line.width = Pt(2)

# Main circle
circle_size = 2.8
main_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(center_x - circle_size / 2), Inches(center_y - circle_size / 2), Inches(circle_size), Inches(circle_size))
main_circle.fill.solid(); main_circle.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); main_circle.line.fill.background()

# Center text
cb = slide.shapes.add_textbox(Inches(center_x - 1.75), Inches(center_y - 0.55), Inches(3.5), Inches(0.9))
p = cb.text_frame.paragraphs[0]; p.text = "The Loop"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_BG); p.alignment = PP_ALIGN.CENTER

sub = slide.shapes.add_textbox(Inches(center_x - 2), Inches(center_y + 0.3), Inches(4), Inches(0.4))
p = sub.text_frame.paragraphs[0]; p.text = "Virtuous cycle"
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(13); p.font.color.rgb = hex_to_rgb(BRAND_BG); p.alignment = PP_ALIGN.CENTER

# 4 surrounding items with arrows
items = [
    ("01 \u00B7 Survey", "Residents complete an anonymous survey"),
    ("02 \u00B7 Credit", "They earn credit toward a report"),
    ("03 \u00B7 Report", "Anyone can buy a Building Insight Report"),
    ("04 \u00B7 Sinking Fund", "Revenue funds the building\u2019s sinking fund"),
]
radius = 3.0

for i, (item_title, item_desc) in enumerate(items):
    angle = (2 * math.pi * i / 4) - math.pi / 2  # top, right, bottom, left
    item_x = center_x + radius * math.cos(angle)
    item_y = center_y + radius * math.sin(angle)

    # Dot on ring
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(item_x - 0.15), Inches(item_y - 0.15), Inches(0.3), Inches(0.3))
    dot.fill.solid(); dot.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_TERTIARY); dot.line.fill.background()

    # Position label card based on quadrant
    box_w = 3.0
    box_h = 1.0
    if i == 0:  # top
        bx = item_x - box_w / 2; by = item_y - box_h - 0.3
        align = PP_ALIGN.CENTER
    elif i == 1:  # right
        bx = item_x + 0.3; by = item_y - box_h / 2
        align = PP_ALIGN.LEFT
    elif i == 2:  # bottom
        bx = item_x - box_w / 2; by = item_y + 0.3
        align = PP_ALIGN.CENTER
    else:  # left
        bx = item_x - box_w - 0.3; by = item_y - box_h / 2
        align = PP_ALIGN.RIGHT

    # Title
    ib = slide.shapes.add_textbox(Inches(bx), Inches(by), Inches(box_w), Inches(0.45))
    p = ib.text_frame.paragraphs[0]; p.text = item_title
    p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_ACCENT); p.alignment = align
    # Description
    db = slide.shapes.add_textbox(Inches(bx), Inches(by + 0.45), Inches(box_w), Inches(0.55))
    tf = db.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = item_desc
    p.font.name = BRAND_BODY_FONT; p.font.size = Pt(11); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY); p.alignment = align

# === Slide 13: Quote ===
slide = new_slide()

quote_mark = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(2.0), Inches(2.0))
p = quote_mark.text_frame.paragraphs[0]; p.text = "\u201C"
p.font.name = "Georgia"; p.font.size = Pt(200); p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.8), Inches(0.12), Inches(2.5))
accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); accent_bar.line.fill.background()

qb = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(10.5), Inches(2.5))
tf = qb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Residents are rewarded for participating, buyers get the data they need, and every dollar from reports goes straight back into maintaining the buildings we call home."
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(28); p.font.italic = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

attr = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(10.5), Inches(0.6))
p = attr.text_frame.paragraphs[0]; p.text = "\u2014 The BlockScore Loop, operating principle"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(18); p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

# === Slide 14: Section break 04 Subscription Model ===
section_break("04", "The Subscription Model", "Five tiers. One dashboard product. Scales with your portfolio.")

# === Slide 15: Multi-card - 5 tier ladder ===
slide = new_slide()

tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
p = tb.text_frame.paragraphs[0]; p.text = "A five-tier subscription ladder"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.25), Inches(2.5), Inches(0.06))
underline.fill.solid(); underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); underline.line.fill.background()

stb = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.333), Inches(0.5))
p = stb.text_frame.paragraphs[0]; p.text = "Start free with Scout. Move up as your portfolio grows. 30-day free trial on every paid plan."
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(14); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

tiers = [
    ("Scout", "Free", "Self-managed body corps", "1 building \u00B7 5 units"),
    ("Tenant", "From $90/mo", "Single-building body corps", "1 building \u00B7 10 units"),
    ("Manager", "From $149/mo", "Strata managers (popular)", "3 buildings \u00B7 20 units"),
    ("Landlord", "From $249/mo", "Portfolio owners", "Unlimited \u00B7 50 units"),
    ("Mayor", "From $330/mo", "Enterprise managers", "Unlimited \u00B7 199 units"),
]

num_cards = 5
margin = 0.4
gap = 0.2
total_width = 13.333 - (2 * margin) - ((num_cards - 1) * gap)
card_width = total_width / num_cards
card_height = 4.5
start_y = 2.2

# Staircase offset so each card is slightly higher than previous? Actually subscription ladder -> ascending so each goes up. But slides have downward y. Let's do each card slightly LOWER start for ascending feel. Reverse: decrement y with i.
for i, (name, price, aud, cap) in enumerate(tiers):
    x = margin + (i * (card_width + gap))
    y_offset = -i * 0.15  # each tier sits slightly higher (reads as climbing)
    y = start_y + y_offset
    is_highlight = (i == 2)  # Manager
    accent = BRAND_ACCENT_TERTIARY if is_highlight else (BRAND_ACCENT if i % 2 == 0 else BRAND_ACCENT_SECONDARY)
    bg = BRAND_CARD_BG_ALT if is_highlight else BRAND_CARD_BG

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(card_height))
    card.fill.solid(); card.fill.fore_color.rgb = hex_to_rgb(bg)
    if is_highlight:
        card.line.color.rgb = hex_to_rgb(BRAND_ACCENT_TERTIARY)
        card.line.width = Pt(2)
    else:
        card.line.fill.background()

    # Top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(0.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = hex_to_rgb(accent); bar.line.fill.background()

    # Tier number/label
    tlb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.3), Inches(card_width - 0.4), Inches(0.35))
    p = tlb.text_frame.paragraphs[0]; p.text = f"TIER {i + 1}"
    p.font.name = BRAND_BODY_FONT; p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = hex_to_rgb(accent)

    # Name
    nb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.65), Inches(card_width - 0.4), Inches(0.6))
    tf = nb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = name
    p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Audience
    ab = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.4), Inches(card_width - 0.4), Inches(1.0))
    tf = ab.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = aud
    p.font.name = BRAND_BODY_FONT; p.font.size = Pt(11); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    # Cap
    cb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 2.6), Inches(card_width - 0.4), Inches(0.5))
    tf = cb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = cap
    p.font.name = BRAND_BODY_FONT; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    # Price
    pb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 3.6), Inches(card_width - 0.4), Inches(0.5))
    tf = pb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = price
    p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = hex_to_rgb(accent)

    # "Most popular" ribbon
    if is_highlight:
        ribbon = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.1), Inches(y - 0.35), Inches(card_width - 0.2), Inches(0.32))
        ribbon.fill.solid(); ribbon.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_TERTIARY); ribbon.line.fill.background()
        rb = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y - 0.38), Inches(card_width - 0.2), Inches(0.35))
        p = rb.text_frame.paragraphs[0]; p.text = "MOST POPULAR"
        p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(10); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_BG); p.alignment = PP_ALIGN.CENTER

out = Path("C:/dev/obsidian/second-brain-skills/.claude/skills/pptx-generator/output/blockscore/livability-index-2026-04-18-part3.pptx")
prs.save(out)
print(f"Saved {out}")
