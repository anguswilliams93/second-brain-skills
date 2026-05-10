from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BRAND_BG = "0d2e28"
BRAND_BG_ALT = "0a231e"
BRAND_TEXT = "ffffff"
BRAND_TEXT_SECONDARY = "a1e8d9"
BRAND_ACCENT = "4dd4c6"
BRAND_ACCENT_SECONDARY = "2dd4aa"
BRAND_ACCENT_TERTIARY = "6ffce0"
BRAND_CARD_BG = "113a33"
BRAND_CARD_BG_ALT = "15473e"
BRAND_HEADING_FONT = "Inter Tight"
BRAND_BODY_FONT = "Inter"


def hex_to_rgb(h):
    h = h.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = hex_to_rgb(BRAND_BG)
    return s


def multi_card_slide(title, cards, subtitle=None):
    slide = new_slide()
    # Title
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.9))
    p = tb.text_frame.paragraphs[0]; p.text = title
    p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(36); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

    underline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35), Inches(2.5), Inches(0.06))
    underline.fill.solid(); underline.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); underline.line.fill.background()

    if subtitle:
        stb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5))
        p = stb.text_frame.paragraphs[0]; p.text = subtitle
        p.font.name = BRAND_BODY_FONT; p.font.size = Pt(15); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

    num_cards = len(cards)
    margin = 0.5
    gap = 0.3
    total_width = 13.333 - (2 * margin) - ((num_cards - 1) * gap)
    card_width = total_width / num_cards
    card_height = 4.6
    start_y = 2.1 if subtitle else 1.8

    accent_colors = [BRAND_ACCENT, BRAND_ACCENT_SECONDARY, BRAND_ACCENT_TERTIARY, BRAND_ACCENT, BRAND_ACCENT_SECONDARY]

    for i, (ctitle, cdesc) in enumerate(cards):
        x = margin + (i * (card_width + gap))
        y = start_y
        accent = accent_colors[i % len(accent_colors)]

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(card_height))
        card.fill.solid(); card.fill.fore_color.rgb = hex_to_rgb(BRAND_CARD_BG); card.line.fill.background()

        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(card_width), Inches(0.1))
        bar.fill.solid(); bar.fill.fore_color.rgb = hex_to_rgb(accent); bar.line.fill.background()

        nb = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.3), Inches(0.8), Inches(0.6))
        pp = nb.text_frame.paragraphs[0]; pp.text = f"0{i + 1}"
        pp.font.name = BRAND_HEADING_FONT; pp.font.size = Pt(24); pp.font.bold = True; pp.font.color.rgb = hex_to_rgb(accent)

        tt = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.0), Inches(card_width - 0.4), Inches(1.2))
        tff = tt.text_frame; tff.word_wrap = True
        pp = tff.paragraphs[0]; pp.text = ctitle
        pp.font.name = BRAND_HEADING_FONT; pp.font.size = Pt(17); pp.font.bold = True; pp.font.color.rgb = hex_to_rgb(BRAND_TEXT)

        dd = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 2.2), Inches(card_width - 0.4), Inches(2.2))
        tff = dd.text_frame; tff.word_wrap = True
        pp = tff.paragraphs[0]; pp.text = cdesc
        pp.font.name = BRAND_BODY_FONT; pp.font.size = Pt(12); pp.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)


# === Slide 6: Benchmark bands (multi-card, 5 cards) ===
multi_card_slide(
    "Benchmark bands",
    [
        ("Poor", "0\u201349 \u2014 well below national wellbeing norms. Immediate attention needed."),
        ("Below avg", "50\u201364 \u2014 lagging Australian wellbeing norms. Clear room to move."),
        ("Average", "65\u201374 \u2014 within normal national range. Stable but unremarkable."),
        ("Good", "75\u201384 \u2014 at or above the Australian wellbeing norm (~75)."),
        ("Excellent", "85\u2013100 \u2014 top decile. Rare and reputationally significant."),
    ],
    subtitle="Five interpretation bands, calibrated against Australian wellbeing research",
)

# === Slide 7: Continuous ML (multi-card, 5 cards) ===
multi_card_slide(
    "Continuously refined with ML",
    [
        ("Collect", "Anonymous responses flow in continuously from tenants across every building on the platform."),
        ("Validate", "Psychometric checks confirm the instrument still measures what it should \u2014 building and portfolio scale."),
        ("Calibrate", "Layer weights are periodically re-fit against held-out behavioural outcomes so the score stays predictive."),
        ("Detect", "Trimmed statistics and confidence ribbons dampen the impact of over-negative or over-positive voices."),
        ("Refine", "Index revisions are version-controlled. Models improve; historical benchmarks stay interpretable."),
    ],
    subtitle="The Index is a living model, not a static questionnaire",
)

# === Slide 8: Section break 02 Products ===
slide = new_slide()
accent_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.5), Inches(7.5))
accent_block.fill.solid(); accent_block.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT); accent_block.line.fill.background()

num = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(3.5), Inches(1.5))
p = num.text_frame.paragraphs[0]; p.text = "02"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(120); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_ACCENT); p.alignment = PP_ALIGN.CENTER

divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(1.5), Inches(0.08), Inches(4.5))
divider.fill.solid(); divider.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); divider.line.fill.background()

tb = slide.shapes.add_textbox(Inches(5.2), Inches(2.6), Inches(7.5), Inches(1.5))
p = tb.text_frame.paragraphs[0]; p.text = "The Products"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

sb2 = slide.shapes.add_textbox(Inches(5.2), Inches(4.2), Inches(7.5), Inches(1.0))
p = sb2.text_frame.paragraphs[0]; p.text = "Three products pointing at the same index from different angles."
p.font.name = BRAND_BODY_FONT; p.font.size = Pt(20); p.font.color.rgb = hex_to_rgb(BRAND_TEXT_SECONDARY)

# === Slide 9: Products (multi-card, 3 cards) ===
multi_card_slide(
    "Three products, one index",
    [
        ("Survey App", "Residents answer the Livability Index anonymously on mobile and earn credit toward a paid report."),
        ("Dashboard", "Live Livability Index for body corporates and strata managers. Layer breakdowns, trends, risk alerts, benchmarks."),
        ("Reports", "Property-specific Livability reports for buyers, renters and owners. Revenue flows into the sinking fund."),
    ],
    subtitle="Residents generate the evidence, managers act on it, buyers pay for it.",
)

# === Slide 10: Two-column Dashboard vs Reports ===
slide = new_slide()

title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
title_bar.fill.solid(); title_bar.fill.fore_color.rgb = hex_to_rgb(BRAND_BG_ALT); title_bar.line.fill.background()

tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
p = tb.text_frame.paragraphs[0]; p.text = "Dashboard vs Reports"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(32); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

# Center divider
divider = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.63), Inches(1.5), Inches(0.07), Inches(5.5))
divider.fill.solid(); divider.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); divider.line.fill.background()

# LEFT (Dashboard)
left_ind = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.7), Inches(0.15), Inches(0.5))
left_ind.fill.solid(); left_ind.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); left_ind.line.fill.background()

lh = slide.shapes.add_textbox(Inches(0.85), Inches(1.65), Inches(5.5), Inches(0.6))
p = lh.text_frame.paragraphs[0]; p.text = "Dashboard \u2014 SaaS subscription"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_ACCENT)

LEFT_BULLETS = [
    "Recurring monthly or annual pricing",
    "Live scores, trends and benchmarks",
    "Multi-building portfolio views",
    "Risk alerts and AI-powered insights",
]
for i, bullet in enumerate(LEFT_BULLETS):
    y_pos = 2.5 + (i * 1.0)
    marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y_pos + 0.12), Inches(0.1), Inches(0.1))
    marker.fill.solid(); marker.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT); marker.line.fill.background()
    txt = slide.shapes.add_textbox(Inches(0.85), Inches(y_pos), Inches(5.5), Inches(0.9))
    tf = txt.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = bullet
    p.font.name = BRAND_BODY_FONT; p.font.size = Pt(18); p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

# RIGHT (Reports)
right_ind = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.93), Inches(1.7), Inches(0.15), Inches(0.5))
right_ind.fill.solid(); right_ind.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_TERTIARY); right_ind.line.fill.background()

rh = slide.shapes.add_textbox(Inches(7.28), Inches(1.65), Inches(5.5), Inches(0.6))
p = rh.text_frame.paragraphs[0]; p.text = "Reports \u2014 One-off purchase"
p.font.name = BRAND_HEADING_FONT; p.font.size = Pt(22); p.font.bold = True; p.font.color.rgb = hex_to_rgb(BRAND_ACCENT_TERTIARY)

RIGHT_BULLETS = [
    "Pay once, per building",
    "Buyer, renter and owner audiences",
    "Revenue funds the building\u2019s sinking fund",
    "Residents redeem survey credit against cost",
]
for i, bullet in enumerate(RIGHT_BULLETS):
    y_pos = 2.5 + (i * 1.0)
    marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.93), Inches(y_pos + 0.12), Inches(0.1), Inches(0.1))
    marker.fill.solid(); marker.fill.fore_color.rgb = hex_to_rgb(BRAND_ACCENT_TERTIARY); marker.line.fill.background()
    txt = slide.shapes.add_textbox(Inches(7.28), Inches(y_pos), Inches(5.5), Inches(0.9))
    tf = txt.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = bullet
    p.font.name = BRAND_BODY_FONT; p.font.size = Pt(18); p.font.color.rgb = hex_to_rgb(BRAND_TEXT)

out = Path("C:/dev/obsidian/second-brain-skills/.claude/skills/pptx-generator/output/blockscore/livability-index-2026-04-18-part2.pptx")
prs.save(out)
print(f"Saved {out}")
